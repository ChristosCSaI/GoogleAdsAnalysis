from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Google Ads Performance Analysis"
subtitle.text = "Summary and Recommendations"

# Summary Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Overall Performance Summary"
content.text = (
    "Summarized the overall performance metrics over the two-month period, including:\n"
    "- Clicks\n"
    "- Impressions\n"
    "- CTR\n"
    "- Cost\n"
    "- Average Position\n"
    "- Conversions\n"
    "- Cost per Conversion\n"
    "- Conversion Rate\n"
    "- View-through Conversions\n"
    "- Total Conversion Value"
)

# Visualizations: Clicks, Impressions, and Cost over Time
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Daily Clicks, Impressions, and Cost"
slide.shapes.add_picture('daily_metrics.png', Inches(0.5), Inches(1.5), width=Inches(9))

# Add descriptive text for daily metrics
description = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
description_text_frame = description.text_frame
description_text_frame.text = (
    "The daily trends show fluctuations in clicks, impressions, and cost over the two-month period. "
    "Significant spikes or drops in these metrics can indicate the impact of specific campaigns or external factors."
)

# Add campaign images to slides with descriptions
descriptions = [
    "Campaign 1 has the highest number of clicks and conversions with a relatively low cost per conversion.",
    "Campaign 15 has a very high number of impressions but a low CTR, indicating that while the ads are being displayed frequently, they are not being clicked as often.",
    "Campaign 7 has a high cost per conversion, suggesting it might be less efficient compared to other campaigns.",
    "Conversions are highest for Campaign 1, indicating effective performance.",
    "Campaign 1 also leads in CTR, showing its effectiveness in attracting clicks.",
    "Conversion rate analysis indicates that Campaign 1 is performing well in converting clicks into actions."
]

campaign_images = [
    'campaign_clicks.png',
    'campaign_impressions.png',
    'campaign_cost.png',
    'campaign_conversions.png',
    'campaign_ctr.png',
    'campaign_conversion_rate.png'
]

for img_path, desc in zip(campaign_images, descriptions):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = img_path.split('.')[0].replace('_', ' ').title()
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    description = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    description_text_frame = description.text_frame
    description_text_frame.text = desc

# Add ad group images to slides with descriptions
adgroup_descriptions = [
    "Ad Group 104 has a high number of clicks and conversions with a good CTR.",
    "Ad Group 536 has a very high number of clicks and impressions but a low conversion rate, indicating inefficiency.",
    "Ad Group 103 has a high CTR and conversion rate, indicating good performance.",
    "Conversions are highest for Ad Group 104, indicating effective performance.",
    "Ad Group 104 leads in CTR, showing its effectiveness in attracting clicks.",
    "Conversion rate analysis indicates that Ad Group 103 is performing well in converting clicks into actions."
]

adgroup_images = [
    'adgroup_clicks.png',
    'adgroup_impressions.png',
    'adgroup_cost.png',
    'adgroup_conversions.png',
    'adgroup_ctr.png',
    'adgroup_conversion_rate.png'
]

for img_path, desc in zip(adgroup_images, adgroup_descriptions):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = img_path.split('.')[0].replace('_', ' ').title()
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    description = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    description_text_frame = description.text_frame
    description_text_frame.text = desc

# Save the updated presentation with fixed charts and descriptions
pptx_file_updated = r'C:\Users\user\GoogleAdsPerformance\Google_Ads_Performance_Analysis_with_Descriptions.pptx'
prs.save(pptx_file_updated)

print(f"Presentation saved as {pptx_file_updated}")
