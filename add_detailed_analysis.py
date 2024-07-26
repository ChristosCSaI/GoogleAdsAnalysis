from pptx import Presentation
from pptx.util import Inches

# Load the existing PowerPoint presentation
pptx_file = 'C:/Users/user/GoogleAdsPerformance/Google_Ads_Performance_Analysis_with_Recommendations_final.pptx'
prs = Presentation(pptx_file)

# Create a new slide for detailed analysis of what happened
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Detailed Analysis of Performance"

content.text = (
    "Based on the visualizations, the following insights were observed:\n\n"
    "1. **Overall Performance Summary:**\n"
    "- The overall performance metrics provide a broad view of the campaign effectiveness over the two-month period, including clicks, impressions, CTR, cost, and conversions.\n\n"
    "2. **Daily Clicks, Impressions, and Cost:**\n"
    "- Fluctuations in daily metrics indicate the impact of specific campaigns or external factors. Significant spikes or drops suggest variations in user engagement and ad spending.\n\n"
    "3. **Total Clicks by Campaign:**\n"
    "- Campaign 1 has the highest number of clicks and conversions with a relatively low cost per conversion, indicating high effectiveness.\n\n"
    "4. **Total Impressions by Campaign:**\n"
    "- Campaign 15 shows high impressions but low CTR, suggesting the need for better ad copy or targeting to improve user engagement.\n\n"
    "5. **Total Cost by Campaign:**\n"
    "- Campaign 7 has a high cost per conversion, indicating inefficiency and the need for cost optimization strategies.\n\n"
    "6. **Total Conversions by Campaign:**\n"
    "- Campaign 1 leads in conversions, demonstrating effective performance in driving desired user actions.\n\n"
    "7. **CTR by Campaign:**\n"
    "- Campaign 1 also leads in CTR, showing its effectiveness in attracting clicks through well-optimized ads.\n\n"
    "8. **Conversion Rate by Campaign:**\n"
    "- Conversion rate analysis indicates that Campaign 1 is performing well in converting clicks into actions.\n\n"
    "9. **Top 10 Ad Groups by Clicks:**\n"
    "- Ad Group 104 has high clicks and conversions with a good CTR, indicating strong performance.\n\n"
    "10. **Top 10 Ad Groups by Impressions:**\n"
    "- Ad Group 536 shows high clicks and impressions but low conversion rate, indicating inefficiency and the need for optimization.\n\n"
    "11. **Top 10 Ad Groups by Cost:**\n"
    "- Ad Group 103 has a high CTR and conversion rate, indicating good performance.\n\n"
    "12. **Top 10 Ad Groups by Conversions:**\n"
    "- Ad Group 104 leads in conversions, showing effective performance.\n\n"
    "13. **Top 10 Ad Groups by CTR:**\n"
    "- Ad Group 104 shows high CTR, demonstrating its effectiveness in attracting clicks.\n\n"
    "14. **Top 10 Ad Groups by Conversion Rate:**\n"
    "- Conversion rate analysis indicates that Ad Group 103 is performing well in converting clicks into actions."
)

# Save the updated presentation with detailed analysis
pptx_file_updated_with_analysis = 'C:/Users/user/GoogleAdsPerformance/Google_Ads_Performance_Analysis_with_Detailed_Analysis.pptx'
prs.save(pptx_file_updated_with_analysis)

print(f"Presentation saved as {pptx_file_updated_with_analysis}")
