from pptx import Presentation
from pptx.util import Inches

# Initialize the presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Google Ads Performance Analysis"
subtitle.text = "Overview of Campaign Performance and Recommendations for Improvement"

# Add a slide with key insights
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Insights"
content.text = ("1. Focus on campaigns with high impressions but low clicks to improve CTR.\n"
                "2. Analyze and optimize campaigns with low CTR or high CPC.\n"
                "3. Invest in high-converting campaigns and identify success factors.\n"
                "4. Reduce cost per conversion by optimizing or pausing high-cost campaigns.\n"
                "5. Monitor and replicate strategies from high conversion rate campaigns.")

# Add visualizations
# Function to add image to a slide
def add_image_to_slide(prs, title_text, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    img_path = image_path
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Save plots as images and add to presentation
fig, axs = plt.subplots(1, 1, figsize=(10, 6))
axs.bar(campaign_summary['Campaign'], campaign_summary['Clicks'], label='Clicks', alpha=0.6)
axs.bar(campaign_summary['Campaign'], campaign_summary['Impressions'], label='Impressions', alpha=0.6)
axs.set_title('Clicks and Impressions by Campaign')
axs.set_xlabel('Campaign')
axs.set_ylabel('Count')
axs.legend()
image_path = "/mnt/data/clicks_impressions.png"
plt.savefig(image_path)
add_image_to_slide(prs, 'Clicks and Impressions by Campaign', image_path)
plt.close()

fig, axs = plt.subplots(1, 1, figsize=(10, 6))
axs.scatter(campaign_summary['CTR'], campaign_summary['Avg_CPC'], color='purple')
axs.set_title('CTR and Avg. CPC by Campaign')
axs.set_xlabel('CTR')
axs.set_ylabel('Avg. CPC')
image_path = "/mnt/data/ctr_avg_cpc.png"
plt.savefig(image_path)
add_image_to_slide(prs, 'CTR and Avg. CPC by Campaign', image_path)
plt.close()

fig, axs = plt.subplots(1, 1, figsize=(10, 6))
axs.bar(campaign_summary['Campaign'], campaign_summary['Conversions'], color='green')
axs.set_title('Conversions by Campaign')
axs.set_xlabel('Campaign')
axs.set_ylabel('Conversions')
image_path = "/mnt/data/conversions.png"
plt.savefig(image_path)
add_image_to_slide(prs, 'Conversions by Campaign', image_path)
plt.close()

fig, axs = plt.subplots(1, 1, figsize=(10, 6))
axs.bar(campaign_summary['Campaign'], campaign_summary['Cost_per_conversion'], color='orange')
axs.set_title('Cost per Conversion by Campaign')
axs.set_xlabel('Campaign')
axs.set_ylabel('Cost per Conversion')
image_path = "/mnt/data/cost_per_conversion.png"
plt.savefig(image_path)
add_image_to_slide(prs, 'Cost per Conversion by Campaign', image_path)
plt.close()

fig, axs = plt.subplots(1, 1, figsize=(10, 6))
axs.plot(campaign_summary['Campaign'], campaign_summary['Conversion_rate'], marker='o', linestyle='-', color='red')
axs.set_title('Conversion Rate by Campaign')
axs.set_xlabel('Campaign')
axs.set_ylabel('Conversion Rate')
image_path = "/mnt/data/conversion_rate.png"
plt.savefig(image_path)
add_image_to_slide(prs, 'Conversion Rate by Campaign', image_path)
plt.close()

# Save the presentation
presentation_path = "/mnt/data/Google_Ads_Performance_Analysis.pptx"
prs.save(presentation_path)

presentation_path
