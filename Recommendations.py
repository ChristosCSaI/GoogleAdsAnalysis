from pptx import Presentation
from pptx.util import Inches

# Load the existing PowerPoint presentation
pptx_file = 'C:/Users/user/GoogleAdsPerformance/Google_Ads_Performance_Analysis_with_Descriptions.pptx'
prs = Presentation(pptx_file)

# Create a new slide for performance improvement recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Performance Improvement Recommendations"

content.text = (
    "Based on the analysis of the Google Ads performance data, the following recommendations can help improve campaign performance:\n\n"
    "1. **Optimize Keywords:**\n"
    "- Identify and pause/remove non-performing keywords with low CTR or high CPC.\n"
    "- Example: Review and optimize keywords in Campaign 15 with high impressions but low CTR.\n\n"
    "2. **Refine Ad Copy:**\n"
    "- A/B test different ad copies to improve CTR and conversion rate.\n"
    "- Example: Rewrite ad copy for Campaign 15 to be more compelling.\n\n"
    "3. **Adjust Bids:**\n"
    "- Increase bids on high-performing keywords and decrease bids on underperforming ones.\n"
    "- Example: Increase bids for high-converting keywords in Campaign 1.\n\n"
    "4. **Improve Landing Pages:**\n"
    "- Ensure landing pages are relevant, load quickly, and provide a good user experience.\n"
    "- Example: Improve landing pages for Ad Group 536 to increase conversion rate.\n\n"
    "5. **Use Negative Keywords:**\n"
    "- Add negative keywords to prevent ads from showing for irrelevant searches.\n\n"
    "6. **Schedule Ads:**\n"
    "- Analyze performance by time of day and schedule ads during peak performance times.\n\n"
    "7. **Geographic Targeting:**\n"
    "- Target locations that perform better and exclude non-performing locations.\n\n"
    "8. **Use Ad Extensions:**\n"
    "- Utilize all relevant ad extensions to improve visibility and CTR.\n\n"
    "9. **Optimize for Mobile:**\n"
    "- Ensure ads and landing pages are mobile-friendly to increase conversion rates.\n\n"
    "10. **Budget Allocation:**\n"
    "- Allocate more budget to high-performing campaigns and ad groups.\n"
    "- Example: Increase budget for Campaign 1 with high conversion rate and ROI."
)

# Save the updated presentation
pptx_file_updated = 'C:/Users/user/GoogleAdsPerformance/Google_Ads_Performance_Analysis_with_Recommendations_final.pptx'
prs.save(pptx_file_updated)

print(f"Presentation saved as {pptx_file_updated}")
